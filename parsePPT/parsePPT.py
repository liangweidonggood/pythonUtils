import os
import sys

from pptx import Presentation
from pptx.shapes.picture import Picture

path = os.path.abspath(os.path.dirname(sys.argv[0]))

fileName = "20220613"
dirPath = "./images/" + fileName
isExists = os.path.exists(dirPath)
if not isExists:
    os.makedirs(dirPath)
else:
    print("已经存在目录")
ppt = Presentation(path + "\\" + fileName + ".pptx")

index = 1

code = ""
name = ""
specs = ""
years = ""
culture = ""
origin = ""
remark = ""
# urlPrefix = "https://sunsah.oss-cn-beijing.aliyuncs.com/202110/"
urlPrefix = "https://api.sunsgallery.cn:8443/images/" + fileName + "/"
url = ""
sql = """
INSERT INTO `ah`.`t_gem` ( `code`, `name`, `specs`, `years`, `culture`, `origin`, `remark`, `url` )
VALUES 
"""
values = ""
count = len(ppt.slides)
print("列表数:", count)
i = 0
for slide in ppt.slides:  # > .slides 得到一个列表，包含每个列表slide
    # print(slide)
    for shape in slide.shapes:  # > slide.shapes 形状
        # print(shape)
        if shape.has_text_frame:  # shape.has_text_frame 判断是否有文字
            text_frame = shape.text_frame  # shape.text_frame 获取文字框
            # print(1)
            # print(text_frame.text)
            for paragraph in text_frame.paragraphs:  # text_frame.paragraphs 获取段落
                # print(1)
                # print(paragraph.text)
                # print(re.findall(r"编号：(.+?)", paragraph.text))
                # print(paragraph.text[3:])
                s = paragraph.text
                if "编号：" in s:
                    code = s[len("编号：") :].strip()
                    # print(code)
                elif "名称：" in s:
                    name = s[len("名称：") :].strip()
                    # print(name)
                elif "规格：" in s:
                    specs = s[len("规格：") :].strip()
                    # print(specs)
                elif "年代：" in s:
                    years = s[len("年代：") :].strip()
                    # print(years)
                elif "文化：" in s:
                    culture = s[len("文化：") :].strip()
                    # print(culture)
                elif "分布出处：" in s:
                    origin = s[len("分布出处：") :].strip()
                    # print(origin)
                elif "备注：" in s:
                    remark = s[len("备注：") :].strip()
                    # print(remark)
            # print(index)
            index += 1
        elif isinstance(shape, Picture):  # 有图片
            # shape.image.blob:二进制图像字节流,写入图像文件
            with open(f"images/{fileName}/{code}.jpg", "wb") as f:
                f.write(shape.image.blob)
    url = urlPrefix + code + ".jpg"
    i += 1
    if count == i:
        values = (
            "('"
            + code
            + "','"
            + name
            + "','"
            + specs
            + "','"
            + years
            + "','"
            + culture
            + "','"
            + origin
            + "','"
            + remark
            + "','"
            + url
            + "');"
        )
    else:
        values = (
            "('"
            + code
            + "','"
            + name
            + "','"
            + specs
            + "','"
            + years
            + "','"
            + culture
            + "','"
            + origin
            + "','"
            + remark
            + "','"
            + url
            + "'),"
        )
    sql += values
print(sql)
